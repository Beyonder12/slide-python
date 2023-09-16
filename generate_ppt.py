from pptx import Presentation
from pptx.util import Inches

# Create a new presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Introduction to Programming & Software Development"
subtitle.text = "Empowering Workers and SMK Teachers"

# Slide 2: What is Programming?
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
body = slide_2.placeholders[1]
title.text = "What is Programming?"
body.text = "Programming is the process of creating a set of instructions that tell a computer how to perform a task."

# ... Add similar code for other slides ...

# Slide 10: Conclusion & Call to Action
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
body = slide_10.placeholders[1]
title.text = "Conclusion & Call to Action"
body.text = ("Emphasize the importance of embracing the digital age.\n"
             "Encourage participants to embark on this exciting journey.\n"
             "Provide information about subsequent sessions or resources they can dive into.")

# Save the presentation
prs.save('presentation.pptx')

print("Presentation generated successfully!")
