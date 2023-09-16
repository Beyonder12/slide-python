from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def get_input_for_slide(slide_title):
    """Get user input for a slide."""
    print(f"\n{slide_title}")
    title = input("Enter the title for the slide: ")
    content = input("Enter the content for the slide: ")
    return title, content

def format_title_shape(title_shape):
    """Apply formatting to a title."""
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def format_content_shape(content_shape):
    """Apply formatting to content."""
    content_shape.text_frame.paragraphs[0].font.size = Pt(20)
    content_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)


# Create a new presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Introduction to Programming & Software Development"
subtitle.text = "Empowering Workers and SMK Teachers"

# Interactive part to gather slides
number_of_slides = int(input("\nHow many slides do you want to create? "))
for i in range(number_of_slides):
    slide_title = f"Slide {i + 1}"
    title_text, content_text = get_input_for_slide(slide_title)
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title_text
    content_shape.text = content_text

    format_title_shape(title_shape)
    format_content_shape(content_shape)

# Slide with Conclusion (if desired)
add_conclusion = input("\nDo you want to add a conclusion slide? (yes/no): ").lower()
if add_conclusion == 'yes':
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Conclusion & Call to Action"
    body.text = ("Emphasize the importance of embracing the digital age.\n"
                 "Encourage participants to embark on this exciting journey.\n"
                 "Provide information about subsequent sessions or resources they can dive into.")
    format_title_shape(title)
    format_content_shape(body)

# Save the presentation to the Desktop directory
file_name = input("\nEnter the name for the saved presentation (without .pptx): ")
save_path = os.path.expanduser(f"~/Desktop/{file_name}.pptx")
prs.save(save_path)

print(f"\nPresentation saved to {save_path}!")
