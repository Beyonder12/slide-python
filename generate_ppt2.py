from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a new presentation object
prs = Presentation()

# Helper function to format title
def format_title(title):
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Use Title Slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Introduction to Programming & Software Development"
subtitle.text = "Empowering Workers and SMK Teachers"
format_title(title)

# Slides 2 to 9: Content slides
titles = [
    "What is Programming?",
    "History of Programming",
    "Programming Languages",
    "Software Development Life Cycle",
    "Modern Development Tools",
    "Introduction to Algorithms",
    "The Role of Debugging",
    "Future Trends in Software"
]

contents = [
    "Programming is the process of creating a set of instructions for a computer.",
    "From punch cards to modern IDEs, programming has a rich history.",
    "Different tasks require different languages, from Python to Java and beyond.",
    "From requirements to deployment and maintenance.",
    "Integrated Development Environments and version control.",
    "Algorithms are sets of rules to solve problems.",
    "Debugging is crucial for ensuring code runs correctly.",
    "Cloud computing, AI in coding, and more exciting futures ahead."
]

for i in range(len(titles)):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use Title and Content layout
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = titles[i]
    content.text = contents[i]
    format_title(title)
    
    # Adding a decorative shape directly (without relying on placeholders)
    left = Inches(4)
    top = Inches(1.5)
    width = Inches(3)
    height = Inches(3)
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

# Slide 10: Conclusion & Call to Action
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.placeholders[1]
title.text = "Conclusion & Call to Action"
content.text = ("Embrace the digital age. Dive into programming. "
                "Continue learning and expanding your skills!")
format_title(title)

# Save the presentation to the Desktop directory
save_path = os.path.expanduser("~/Desktop/presentation.pptx")
prs.save(save_path)

print(f"Presentation saved to {save_path}!")
